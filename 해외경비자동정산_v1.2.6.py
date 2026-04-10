import os
import queue
import re
import shutil
import sys
import threading
import tkinter as tk
from io import BytesIO
from pathlib import Path
from tkinter import StringVar, Tk, filedialog, messagebox, ttk
from tkinter.scrolledtext import ScrolledText

from PIL import Image, ImageFilter, ImageOps, ImageTk
import pytesseract

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm, Pt

    PPTX_AVAILABLE = True
except Exception:
    Presentation = None
    PP_ALIGN = None
    Cm = None
    Pt = None
    PPTX_AVAILABLE = False

try:
    from openpyxl import load_workbook

    OPENPYXL_AVAILABLE = True
except Exception:
    load_workbook = None
    OPENPYXL_AVAILABLE = False

DEFAULT_TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png"}
UNCATEGORIZED_FOLDER_NAME = "미분류"
EXCEL_BASENAME = "해외출장비정산서_RMB"
EXCEL_EXTENSIONS = (".xlsx", ".xlsm", ".xltx", ".xltm")
EXCEL_SUM_SHEET_NAME = "sum"
EXCEL_SUM_SHEET_CANDIDATES = ("sum", "경비sum", "경비 sum")
EXCEL_HEADER_SCAN_MAX_ROW = 80
EXCEL_HEADER_SCAN_MAX_COL = 80
PPTX_EXTENSIONS = (".pptx",)
APP_BASENAME = "해외경비자동정산"
_runtime_stem = Path(sys.executable if getattr(sys, "frozen", False) else __file__).stem
_version_match = re.search(r"(v\d+\.\d+\.\d+|v\d+\.\d+)", _runtime_stem, re.IGNORECASE)
APP_VERSION = _version_match.group(1) if _version_match else "v1.0.0"
APP_NAME = f"{APP_BASENAME}_{APP_VERSION}"
RUN_MODE_OPTIONS = {
    "1": ("1. 그림인식해서 폴더안에 넣기(PPT는 생성하지 않음)", True, False),
    "2": ("2. 폴더안의 사진 분류하여 PPT만들기", False, True),
    "3": ("3. 위의 2가지 기능 같이 구현하기", True, True),
}
DESC_FILE_NAME = "desc.txt"
DEFAULT_INPUT_FOLDER = r"C:\VAD_PC\경비\해외_출장경비"
APP_ICON_CANDIDATES = ("vad_app.ico", "logo.ico", "logo-small.png", "logo-footer.png", "VAD Logo(XL).jpg")
BACKGROUND_IMAGE_CANDIDATES = ("IMG_1364.JPG", "VAD_banner2.png", "VAD_banner.png", "ceo_01.jpg", "ceo_02.jpg")
BACKGROUND_BLUR_RADIUS = 12
HELP_WINDOW_GEOMETRY = "760x520"
PREVIEW_MIN_SIZE = (280, 280)
PREVIEW_BG_COLOR = "#f6f8fb"
RESAMPLE_LANCZOS = Image.Resampling.LANCZOS if hasattr(Image, "Resampling") else Image.LANCZOS
LOG_PROCESSING_MARKER = "[[PROCESSING]]"
TOTAL_PIC_SAMPLE_PATHS = (
    r"C:\Users\kkich\Desktop\mikedb001\vad_expenses_overseas\total_pic_sample\1.png",
    r"C:\Users\kkich\Desktop\mikedb001\vad_expenses_overseas\total_pic_sample\2.png",
    r"C:\Users\kkich\Desktop\mikedb001\vad_expenses_overseas\total_pic_sample\3.png",
)

FULLWIDTH_TRANS = str.maketrans(
    {
        "０": "0",
        "１": "1",
        "２": "2",
        "３": "3",
        "４": "4",
        "５": "5",
        "６": "6",
        "７": "7",
        "８": "8",
        "９": "9",
        "．": ".",
        "，": ",",
        "￥": "¥",
    }
)

NEGATIVE_AMOUNT_RE = re.compile(r"-\s*([0-9][0-9,]*\.[0-9]{2})")
CURRENCY_AMOUNT_RE = re.compile(r"(?:[¥￥]|RMB|CNY)\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)", re.IGNORECASE)
CONTEXT_AMOUNT_RE = re.compile(
    r"(?:합계|총계|총액|금액|결제금액|실결제|실지불|实付|应付|合计|总计|总额|金额)\D{0,6}([0-9][0-9,]*(?:\.[0-9]{1,2})?)",
    re.IGNORECASE,
)
GENERIC_NUMBER_RE = re.compile(r"([0-9]{1,3}(?:,[0-9]{3})*(?:\.[0-9]{1,2})?|[0-9]+(?:\.[0-9]{1,2})?)")
TOTAL_CONTEXT_AMOUNT_RE = re.compile(
    r"(?:总金额|总计|合计|应付|实付|支付金额|결제금액|총금액|합계|total)\D{0,8}-?\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)",
    re.IGNORECASE,
)
NEGATIVE_OR_PLAIN_AMOUNT_RE = re.compile(r"-\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)")
DASH_AMOUNT_RE = re.compile(r"-\s*([0-9][0-9,]*(?:\.[0-9]{1,2})?)")
SEARCH_TOP_RATIO = 0.01
SEARCH_BOTTOM_RATIO_DARK_BG = 0.40
SEARCH_BOTTOM_RATIO_LIGHT_BG = 0.58
BACKGROUND_DARK_THRESHOLD = 132

CATEGORIES = {
    "교통비": ["滴滴出行", "师傅", "船票", "机票", "火车", "乘车码", "邮轮母港", "客运", "汽车", "出租车", "铁路", "交通卡"],
    "개인경비": ["星巴克", "美团订", "喜茶", "咖啡", "便利", "按摩", "贸易", "超市", "果子", "贡茶", "COFFEE", "麦当劳", "友宝", "免税", "烤串", "美宜佳", "豪士特", "常知乐", "好利来", "小米之家", "小吃店", "锅巴", "博物馆", "乌龙茶", "休息中心", "购物", "商行", "骨科"],
    "숙박비": ["酒店"],
    "식대": ["餐厅", "串", "肯德基", "餐饮", "拉面", "烤肉", "菜馆", "烧烤", "韩餐", "小草", "火锅", "料理", "点餐"],
    "통신비": ["手机", "中国移动"],
}

EXCEL_HEADER_KEYWORDS = {
    "교통비": ["교통비", "교통", "운임", "차량", "택시", "transport"],
    "개인경비": ["개인경비", "개인", "잡비", "기타", "etc"],
    "숙박비": ["숙박비", "숙박", "호텔", "hotel"],
    "식대": ["식대", "식사", "식음", "meal", "food"],
    "통신비": ["통신비", "통신", "유심", "전화", "인터넷", "telecom"],
    UNCATEGORIZED_FOLDER_NAME: ["미분류", "기타", "other"],
}

# Windows 환경에서 Tesseract 경로를 기본값으로 지정합니다.
pytesseract.pytesseract.tesseract_cmd = DEFAULT_TESSERACT_PATH


def get_app_dir():
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent


def find_existing_file(candidates):
    base_dir = get_app_dir()
    search_dirs = (base_dir, base_dir / "pics")
    for directory in search_dirs:
        for name in candidates:
            path = directory / name
            if path.exists():
                return path
    return None


def read_text_with_fallback(path):
    encodings = ("utf-8-sig", "utf-8", "cp949", "euc-kr")
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def normalize_text(text):
    return re.sub(r"\s+", "", text)


def normalize_lookup_text(value):
    return re.sub(r"\s+", "", str(value or "")).lower()


def normalize_match_key(value):
    return re.sub(r"[\W_]+", "", str(value or "")).lower()


def categorize_image_with_match(text):
    for category, keywords in CATEGORIES.items():
        for keyword in keywords:
            if keyword in text:
                return category, keyword
    return None, None


def parse_amount_token(token):
    token = token.replace(",", "").strip()
    if token.count(".") > 1:
        return None

    try:
        value = abs(float(token))
    except (TypeError, ValueError):
        return None

    if value == 0 or value > 1_000_000:
        return None
    return round(value, 2)


def extract_amount_from_text(raw_text):
    if not raw_text:
        return None

    text = str(raw_text).translate(FULLWIDTH_TRANS)
    candidates = []

    # 우선순위 1: "-255.00"처럼 음수 기호가 붙은 금액(절대값으로 저장)
    for match in NEGATIVE_AMOUNT_RE.finditer(text):
        value = parse_amount_token(match.group(1))
        if value is not None:
            candidates.append((4, value))

    for match in CURRENCY_AMOUNT_RE.finditer(text):
        value = parse_amount_token(match.group(1))
        if value is not None:
            candidates.append((3, value))

    for match in CONTEXT_AMOUNT_RE.finditer(text):
        value = parse_amount_token(match.group(1))
        if value is not None:
            candidates.append((2, value))

    for match in GENERIC_NUMBER_RE.finditer(text):
        value = parse_amount_token(match.group(1))
        if value is None:
            continue
        if float(value).is_integer() and 1900 <= int(value) <= 2100:
            continue
        candidates.append((1, value))

    if not candidates:
        return None

    top_score = max(score for score, _ in candidates)
    top_values = [value for score, value in candidates if score == top_score]
    return round(max(top_values), 2)


def categorize_image(text):
    category, _ = categorize_image_with_match(text)
    return category


def collect_images(folder_path):
    folder = Path(folder_path)
    return sorted(
        [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() in IMAGE_EXTENSIONS],
        key=lambda p: p.name.lower(),
    )


def resolve_run_mode(mode_code):
    return RUN_MODE_OPTIONS.get(mode_code, RUN_MODE_OPTIONS["1"])


def process_images(input_folder, log_callback=print, progress_callback=None, preview_callback=None):
    input_path = Path(input_folder)
    image_files = collect_images(input_path)
    total_images = len(image_files)
    summary = {
        "total": 0,
        "moved": 0,
        "unmatched": 0,
        "errors": 0,
        "amount_found": 0,
        "amount_missing": 0,
        "category_amounts": {category: [] for category in CATEGORIES},
        "uncategorized_amounts": [],
        "folder_counts": {},
    }

    if progress_callback:
        progress_callback(0, total_images, None)

    for index, file_path in enumerate(image_files, start=1):
        summary["total"] += 1
        try:
            if preview_callback:
                preview_callback(str(file_path))

            with Image.open(file_path) as img:
                ocr_text = pytesseract.image_to_string(img, lang="chi_sim")
                normalized = normalize_text(ocr_text)

            category, matched_keyword = categorize_image_with_match(normalized)
            if category and matched_keyword:
                log_callback(
                    f"{LOG_PROCESSING_MARKER}[인식중] {file_path.name} | 분류일치: {matched_keyword} -> {category}"
                )
            else:
                log_callback(f"{LOG_PROCESSING_MARKER}[인식중] {file_path.name} | 분류일치 없음")

            amount = extract_amount_from_text(ocr_text)

            if amount is not None:
                summary["amount_found"] += 1
            else:
                summary["amount_missing"] += 1

            if category:
                target_folder = input_path / category
                target_folder.mkdir(parents=True, exist_ok=True)
                shutil.move(str(file_path), str(target_folder / file_path.name))
                summary["moved"] += 1
                summary["folder_counts"][category] = summary["folder_counts"].get(category, 0) + 1
                if amount is not None:
                    summary["category_amounts"][category].append(amount)
                log_callback(f"[처리] {file_path.name} -> {category}")
            else:
                summary["unmatched"] += 1
                unmatched_folder = input_path / UNCATEGORIZED_FOLDER_NAME
                unmatched_folder.mkdir(parents=True, exist_ok=True)
                shutil.move(str(file_path), str(unmatched_folder / file_path.name))
                summary["folder_counts"][UNCATEGORIZED_FOLDER_NAME] = (
                    summary["folder_counts"].get(UNCATEGORIZED_FOLDER_NAME, 0) + 1
                )
                if amount is not None:
                    summary["uncategorized_amounts"].append(amount)
                log_callback(f"[처리] {file_path.name} -> {UNCATEGORIZED_FOLDER_NAME}")
        except Exception as exc:
            summary["errors"] += 1
            log_callback(f"[오류] {file_path.name}: {exc}")
        finally:
            if progress_callback:
                progress_callback(index, total_images, file_path.name)

    summary["category_totals"] = {
        category: round(sum(values), 2) for category, values in summary["category_amounts"].items()
    }
    summary["uncategorized_total"] = round(sum(summary["uncategorized_amounts"]), 2)
    return summary


def create_ppt_from_subfolders(root_folder, output_path=None, log_callback=print):
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx 모듈이 설치되어 있지 않습니다.")

    root_path = Path(root_folder)
    if output_path is None:
        output_path = root_path / f"{root_path.name}.pptx"
    else:
        output_path = Path(output_path)

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    blank_layout = prs.slide_layouts[6]

    photos_per_slide = 16
    pic_width = Cm(4.07)
    pic_height = Cm(7.9)
    left_start = Cm(0.43)
    left_end = Cm(29.37)
    gap = int((left_end - left_start - (pic_width * 7)) / 7)
    row_top = [Cm(1.53), Cm(10.13)]

    title_left = Cm(0.43)
    title_top = Cm(0.23)
    title_width = Cm(14.11)
    title_height = Cm(1.06)

    summary = {
        "folders_total": 0,
        "folders_used": 0,
        "slides_total": 0,
        "images_total": 0,
        "output_path": str(output_path),
        "folder_image_counts": {},
    }

    subfolders = sorted([p for p in root_path.iterdir() if p.is_dir()], key=lambda p: p.name.lower())
    for subfolder in subfolders:
        summary["folders_total"] += 1
        photo_files = collect_images(subfolder)
        if not photo_files:
            continue

        summary["folders_used"] += 1
        summary["folder_image_counts"][subfolder.name] = len(photo_files)

        slide = None
        for i, photo_path in enumerate(photo_files, start=1):
            if ((i - 1) % photos_per_slide) == 0:
                slide = prs.slides.add_slide(blank_layout)
                summary["slides_total"] += 1

                title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
                p = title_shape.text_frame.paragraphs[0]
                p.text = subfolder.name
                p.font.name = "맑은 고딕"
                p.font.size = Pt(18)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT

            idx = (i - 1) % photos_per_slide
            row_idx = idx // 8
            col_idx = idx % 8

            pic_left = left_start + col_idx * (pic_width + gap)
            pic_top = row_top[row_idx]
            slide.shapes.add_picture(str(photo_path), pic_left, pic_top, width=pic_width, height=pic_height)
            summary["images_total"] += 1

    prs.save(str(output_path))
    log_callback(f"[PPT 완료] {output_path}")
    return summary


def find_excel_file(root_folder, base_name=EXCEL_BASENAME):
    root_path = Path(root_folder)
    if not root_path.exists():
        return None

    for ext in EXCEL_EXTENSIONS:
        candidate = root_path / f"{base_name}{ext}"
        if candidate.exists():
            return candidate

    candidates = sorted(
        [
            p
            for p in root_path.iterdir()
            if p.is_file() and p.suffix.lower() in EXCEL_EXTENSIONS and p.stem.startswith(base_name)
        ],
        key=lambda p: p.name.lower(),
    )
    return candidates[0] if candidates else None


def header_matches_category(cell_text, category):
    if not cell_text:
        return False

    norm = normalize_lookup_text(cell_text)
    keys = EXCEL_HEADER_KEYWORDS.get(category, [category])
    for key in keys:
        if normalize_lookup_text(key) in norm:
            return True
    return False


def get_sum_worksheet(workbook):
    targets = {normalize_lookup_text(name) for name in EXCEL_SUM_SHEET_CANDIDATES}
    for worksheet in workbook.worksheets:
        if normalize_lookup_text(worksheet.title) in targets:
            return worksheet
    target_names = ", ".join(EXCEL_SUM_SHEET_CANDIDATES)
    raise RuntimeError(f"엑셀 파일에 '{target_names}' 시트를 찾지 못했습니다.")


def detect_sum_header_row_and_columns(worksheet, categories):
    best_row = None
    best_map = {}

    max_row = min(max(worksheet.max_row, EXCEL_HEADER_SCAN_MAX_ROW), EXCEL_HEADER_SCAN_MAX_ROW)
    max_col = min(max(worksheet.max_column, EXCEL_HEADER_SCAN_MAX_COL), EXCEL_HEADER_SCAN_MAX_COL)

    for row in range(1, max_row + 1):
        row_map = {}
        for col in range(1, max_col + 1):
            value = worksheet.cell(row=row, column=col).value
            if value is None:
                continue
            for category in categories:
                if category in row_map:
                    continue
                if header_matches_category(value, category):
                    row_map[category] = col

        if len(row_map) > len(best_map):
            best_row = row
            best_map = row_map

    if best_row is None or not best_map:
        raise RuntimeError("sum 시트에서 카테고리 헤더 행을 찾지 못했습니다.")

    return best_row, best_map


def find_next_empty_row(worksheet, column, start_row):
    row = max(1, start_row)
    while True:
        cell_value = worksheet.cell(row=row, column=column).value
        if cell_value is None or str(cell_value).strip() == "":
            return row
        row += 1


def detect_sum_header_cells(worksheet):
    best_row = None
    best_headers = {}

    max_row = min(max(worksheet.max_row, EXCEL_HEADER_SCAN_MAX_ROW), EXCEL_HEADER_SCAN_MAX_ROW)
    max_col = min(max(worksheet.max_column, EXCEL_HEADER_SCAN_MAX_COL), EXCEL_HEADER_SCAN_MAX_COL)

    for row in range(1, max_row + 1):
        row_headers = {}
        for col in range(1, max_col + 1):
            value = worksheet.cell(row=row, column=col).value
            if value is None:
                continue
            text = str(value).strip()
            if not text:
                continue
            row_headers[col] = text

        if len(row_headers) > len(best_headers):
            best_row = row
            best_headers = row_headers

    if best_row is None or not best_headers:
        raise RuntimeError("경비sum 시트의 헤더 행을 찾지 못했습니다.")

    return best_row, best_headers


def match_ppt_title_to_header_column(title_text, header_cells):
    title_key = normalize_match_key(title_text)
    if not title_key:
        return None, None

    for col, header_text in header_cells.items():
        header_key = normalize_match_key(header_text)
        if title_key == header_key:
            return col, header_text

    for col, header_text in header_cells.items():
        header_key = normalize_match_key(header_text)
        if not header_key:
            continue
        if title_key in header_key or header_key in title_key:
            return col, header_text

    return None, None


def extract_total_amount_from_text(raw_text, return_score=False):
    if not raw_text:
        if return_score:
            return None, None
        return None

    text = str(raw_text).translate(FULLWIDTH_TRANS)
    candidates = []

    for match in TOTAL_CONTEXT_AMOUNT_RE.finditer(text):
        token = match.group(1)
        value = parse_amount_token(token)
        if value is not None:
            score = 6 if "." in token else 5
            candidates.append((score, value))

    for match in NEGATIVE_OR_PLAIN_AMOUNT_RE.finditer(text):
        token = match.group(1)
        value = parse_amount_token(token)
        if value is not None:
            score = 4 if "." in token else 2
            candidates.append((score, value))

    fallback = extract_amount_from_text(text)
    if fallback is not None:
        candidates.append((3, fallback))

    if not candidates:
        if return_score:
            return None, None
        return None

    top_score = max(score for score, _ in candidates)
    top_values = [value for score, value in candidates if score == top_score]
    selected = round(max(top_values), 2)
    if return_score:
        return selected, top_score
    return selected


def extract_amount_candidates_from_image(image):
    amount_candidates = []
    base = ImageOps.exif_transpose(image).convert("RGB")
    width, height = base.size

    # 샘플(1.png,2.png,3.png) 기준으로 총금액이 하단 중앙에 위치하는 경우가 많아
    # 전체 + 하단 영역을 모두 OCR 대상으로 사용한다.
    regions = [
        base,
        base.crop((0, int(height * 0.45), width, height)),
        base.crop((int(width * 0.2), int(height * 0.38), int(width * 0.8), height)),
    ]

    for region in regions:
        variants = []
        variants.append(region)

        gray = ImageOps.grayscale(region)
        variants.append(gray)

        high_contrast = ImageOps.autocontrast(gray)
        variants.append(high_contrast)

        enlarged = high_contrast.resize(
            (max(1, high_contrast.width * 2), max(1, high_contrast.height * 2)),
            RESAMPLE_LANCZOS,
        )
        variants.append(enlarged)

        binary = enlarged.point(lambda x: 255 if x > 160 else 0, mode="1")
        variants.append(binary)

        for img in variants:
            for psm in ("6", "7"):
                try:
                    ocr_text = pytesseract.image_to_string(img, lang="chi_sim+eng", config=f"--psm {psm}")
                except Exception:
                    continue

                amount, score = extract_total_amount_from_text(ocr_text, return_score=True)
                if amount is not None:
                    amount_candidates.append((amount, score))

    return amount_candidates


def extract_slide_title_text(slide):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        top_pos = int(getattr(shape, "top", 0))
        candidates.append((top_pos, text))

    if not candidates:
        return ""

    candidates.sort(key=lambda x: x[0])
    return candidates[0][1].splitlines()[0].strip()


def extract_slide_total_amount(slide):
    amount_candidates = []

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text_amount, score = extract_total_amount_from_text(shape.text_frame.text, return_score=True)
            if text_amount is not None:
                amount_candidates.append((text_amount, score))

        image_obj = getattr(shape, "image", None)
        if image_obj is None:
            continue

        try:
            with Image.open(BytesIO(image_obj.blob)) as img:
                amount_candidates.extend(extract_amount_candidates_from_image(img))
        except Exception:
            continue

    if not amount_candidates:
        return None

    top_score = max(score for _, score in amount_candidates)
    top_values = [round(value, 2) for value, score in amount_candidates if score == top_score]
    freq_map = {}
    for value in top_values:
        freq_map[value] = freq_map.get(value, 0) + 1

    max_freq = max(freq_map.values())
    frequent_values = [value for value, freq in freq_map.items() if freq == max_freq]

    # OCR 오인식으로 자릿수가 과도하게 커지는 경우를 줄이기 위해
    # 동일 신뢰도/빈도에서는 더 작은 금액을 우선한다.
    return round(min(frequent_values), 2)


def write_ppt_title_matched_amounts_to_excel(pptx_path, excel_path, log_callback=print, progress_callback=None):
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx 모듈이 설치되어 있지 않습니다.")
    if not OPENPYXL_AVAILABLE:
        raise RuntimeError("openpyxl 모듈이 설치되어 있지 않습니다.")

    pptx_path = Path(pptx_path)
    excel_path = Path(excel_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX 파일을 찾을 수 없습니다: {pptx_path}")
    if not excel_path.exists():
        raise FileNotFoundError(f"XLSX 파일을 찾을 수 없습니다: {excel_path}")

    # 총금액 OCR 형태를 샘플과 유사한 패턴으로 추출하기 위한 참고 경로(문서화 목적).
    _ = TOTAL_PIC_SAMPLE_PATHS

    presentation = Presentation(str(pptx_path))
    keep_vba = excel_path.suffix.lower() == ".xlsm"
    workbook = load_workbook(str(excel_path), keep_vba=keep_vba)
    worksheet = get_sum_worksheet(workbook)
    header_row, header_cells = detect_sum_header_cells(worksheet)
    value_row = header_row + 1

    slides_total = len(presentation.slides)
    result = {
        "pptx_path": str(pptx_path),
        "excel_path": str(excel_path),
        "slides_total": slides_total,
        "written_count": 0,
        "error_count": 0,
        "unmatched_titles": [],
        "no_amount_slides": [],
    }

    if progress_callback:
        progress_callback(0, slides_total, None)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        try:
            title_text = extract_slide_title_text(slide)
            log_callback(f"{LOG_PROCESSING_MARKER}[인식중] {pptx_path.name} {slide_index}페이지")

            amount = extract_slide_total_amount(slide)
            if amount is None:
                result["no_amount_slides"].append({"slide_index": slide_index, "title": title_text})
                log_callback(f"[금액없음] {slide_index}페이지 제목 '{title_text or '-'}'")
                if progress_callback:
                    progress_callback(slide_index, slides_total, f"{slide_index}페이지")
                continue

            target_col, header_text = match_ppt_title_to_header_column(title_text, header_cells)
            if target_col is None:
                result["unmatched_titles"].append({"slide_index": slide_index, "title": title_text, "amount": amount})
                log_callback(f"[엑셀 미매핑] {slide_index}페이지 제목 '{title_text or '-'}' -> 일치 헤더 없음")
                if progress_callback:
                    progress_callback(slide_index, slides_total, f"{slide_index}페이지")
                continue

            target_row = find_next_empty_row(worksheet, target_col, value_row)
            worksheet.cell(row=target_row, column=target_col).value = round(amount, 2)
            worksheet.cell(row=target_row, column=target_col).number_format = "#,##0.00"
            result["written_count"] += 1
            log_callback(
                f"[엑셀 입력] {slide_index}페이지 '{title_text or '-'}' -> {header_text}: {amount:.2f} "
                f"(열 {target_col}, 행 {target_row})"
            )
        except Exception as exc:
            result["error_count"] += 1
            log_callback(f"[오류] {slide_index}페이지 처리 실패: {exc}")
        finally:
            if progress_callback:
                progress_callback(slide_index, slides_total, f"{slide_index}페이지")

    try:
        workbook.save(str(excel_path))
    except PermissionError as exc:
        raise RuntimeError("엑셀 파일이 열려 있어 저장할 수 없습니다. 파일을 닫고 다시 실행해 주세요.") from exc

    log_callback(f"[엑셀 완료] 저장: {excel_path}")
    return result


def extract_ppt_titles(pptx_path):
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        return []
    prs = Presentation(str(pptx_path))
    titles = []
    for slide in prs.slides:
        text = extract_slide_title_text(slide).strip()
        if text:
            titles.append(text)
    return titles


def collect_images_by_folder_order(work_folder):
    root = Path(work_folder)
    subfolders = sorted([p for p in root.iterdir() if p.is_dir()], key=lambda p: p.name.lower())

    items = []
    folder_order = []
    if subfolders:
        for folder in subfolders:
            folder_order.append(folder.name)
            for image_path in collect_images(folder):
                items.append((folder.name, image_path))
    else:
        folder_order.append(root.name)
        for image_path in collect_images(root):
            items.append((root.name, image_path))
    return items, folder_order


def build_folder_to_ppt_title_map(folder_order, ppt_titles):
    mapping = {}
    for idx, folder_name in enumerate(folder_order):
        if idx < len(ppt_titles):
            mapping[folder_name] = ppt_titles[idx]
    return mapping


def resolve_header_by_folder(folder_name, header_cells, folder_title_map=None):
    col, header_text = match_ppt_title_to_header_column(folder_name, header_cells)
    if col is not None:
        return col, header_text, f"folder:{folder_name}"

    if folder_title_map:
        mapped_title = folder_title_map.get(folder_name, "")
        if mapped_title:
            col, header_text = match_ppt_title_to_header_column(mapped_title, header_cells)
            if col is not None:
                return col, header_text, f"ppt-title:{mapped_title}"

    return None, None, "unmatched"


def extract_dash_amount_from_text(raw_text, return_score=False):
    if not raw_text:
        if return_score:
            return None, None
        return None

    text = str(raw_text).translate(FULLWIDTH_TRANS)
    text = text.replace("—", "-").replace("–", "-").replace("−", "-")
    candidates = []

    for match in DASH_AMOUNT_RE.finditer(text):
        token = match.group(1)
        value = parse_amount_token(token)
        if value is None:
            continue

        compact = token.replace(",", "")
        dot_count = compact.count(".")
        frac_len = len(compact.split(".")[1]) if "." in compact else 0
        int_len = len(compact.split(".")[0]) if compact else 0

        score = 10.0
        score += 1.8 if dot_count == 1 and frac_len == 2 else (1.0 if dot_count == 1 else 0.4)
        score += min(max(int_len, 1), 8) * 0.06
        candidates.append((score, value))

    if not candidates:
        if return_score:
            return None, None
        return None

    best_score = max(score for score, _ in candidates)
    best_values = [round(value, 2) for score, value in candidates if score == best_score]
    selected = round(max(best_values), 2)
    if return_score:
        return selected, best_score
    return selected


def is_dark_background_image(image):
    gray = ImageOps.grayscale(ImageOps.exif_transpose(image).convert("RGB"))
    width, height = gray.size
    edge_w = max(2, int(width * 0.08))
    top_h = max(2, int(height * 0.22))

    regions = [
        gray.crop((0, 0, width, top_h)),
        gray.crop((0, 0, edge_w, height)),
        gray.crop((max(0, width - edge_w), 0, width, height)),
    ]

    samples = []
    for region in regions:
        samples.extend(region.getdata())

    if not samples:
        return False

    samples.sort()
    median = samples[len(samples) // 2]
    return median < BACKGROUND_DARK_THRESHOLD


def extract_dash_amount_candidates_from_image(image):
    amount_candidates = []
    base = ImageOps.exif_transpose(image).convert("RGB")
    width, height = base.size

    is_dark_bg = is_dark_background_image(base)
    bottom_ratio = SEARCH_BOTTOM_RATIO_DARK_BG if is_dark_bg else SEARCH_BOTTOM_RATIO_LIGHT_BG

    top = max(0, int(height * SEARCH_TOP_RATIO))
    bottom = min(height, max(top + 1, int(height * bottom_ratio)))
    band = base.crop((0, top, width, bottom))
    if band.width < 2 or band.height < 2:
        return amount_candidates

    enlarged = band.resize(
        (max(1, band.width * 3), max(1, band.height * 4)),
        RESAMPLE_LANCZOS,
    )
    gray = ImageOps.grayscale(enlarged)
    auto = ImageOps.autocontrast(gray)
    inv = ImageOps.invert(auto)
    binary_light = auto.point(lambda x: 255 if x > 165 else 0, mode="1")
    binary_dark = inv.point(lambda x: 255 if x > 150 else 0, mode="1")

    if is_dark_bg:
        ocr_targets = (
            (binary_dark, "6"),
            (binary_dark, "7"),
            (inv, "6"),
            (inv, "11"),
        )
    else:
        ocr_targets = (
            (binary_light, "6"),
            (binary_light, "7"),
            (auto, "6"),
            (auto, "11"),
        )

    for img, psm in ocr_targets:
        try:
            cfg = f"--oem 1 --psm {psm} -c tessedit_char_whitelist=0123456789.,-"
            ocr_text = pytesseract.image_to_string(img, lang="chi_sim+eng", config=cfg)
        except Exception:
            continue

        amount, score = extract_dash_amount_from_text(ocr_text, return_score=True)
        if amount is not None:
            amount_candidates.append((amount, score))

    return amount_candidates


def extract_total_amount_from_image(image_path):
    with Image.open(image_path) as img:
        candidates = extract_dash_amount_candidates_from_image(img)

    if not candidates:
        return None

    # 숫자 일치도(동일 값 반복 인식)를 우선하고, 그 다음 OCR 점수를 반영한다.
    stats = {}
    for value, score in candidates:
        key = round(value, 2)
        if key not in stats:
            stats[key] = {"count": 0, "score_sum": 0.0, "score_max": 0.0}
        stats[key]["count"] += 1
        stats[key]["score_sum"] += float(score)
        stats[key]["score_max"] = max(stats[key]["score_max"], float(score))

    best_value = max(
        stats.items(),
        key=lambda kv: (kv[1]["count"], kv[1]["score_max"], kv[1]["score_sum"], kv[0]),
    )[0]
    return round(best_value, 2)


def write_folder_images_to_excel_by_headers(work_folder, excel_path, log_callback=print, progress_callback=None):
    if not OPENPYXL_AVAILABLE:
        raise RuntimeError("openpyxl module is not installed.")

    work_folder = Path(work_folder)
    excel_path = Path(excel_path)

    if not work_folder.exists():
        raise FileNotFoundError(f"Work folder not found: {work_folder}")
    if not excel_path.exists():
        raise FileNotFoundError(f"Excel file not found: {excel_path}")

    image_items, _folder_order = collect_images_by_folder_order(work_folder)

    keep_vba = excel_path.suffix.lower() == ".xlsm"
    workbook = load_workbook(str(excel_path), keep_vba=keep_vba)
    worksheet = get_sum_worksheet(workbook)
    header_row, header_cells = detect_sum_header_cells(worksheet)
    value_row = header_row + 1

    total_images = len(image_items)
    result = {
        "work_folder": str(work_folder),
        "excel_path": str(excel_path),
        "total_images": total_images,
        "written_count": 0,
        "error_count": 0,
        "unmatched_count": 0,
        "no_amount_count": 0,
    }

    if progress_callback:
        progress_callback(0, total_images, None)

    for index, (folder_name, image_path) in enumerate(image_items, start=1):
        try:
            target_col, header_text, class_source = resolve_header_by_folder(folder_name, header_cells)
            amount = extract_total_amount_from_image(image_path)

            amount_text = f"{amount:.2f}" if amount is not None else "없음"
            class_text = header_text if header_text else "미분류"
            log_callback(
                f"{LOG_PROCESSING_MARKER}[인식중] {image_path.name} | 분류: {class_text} ({class_source}) | 금액: {amount_text}"
            )

            if amount is None:
                result["no_amount_count"] += 1
                log_callback(f"[금액없음] {image_path.name}")
                continue

            if target_col is None:
                result["unmatched_count"] += 1
                log_callback(f"[헤더미일치] {image_path.name} -> 폴더 '{folder_name}'")
                continue

            target_row = find_next_empty_row(worksheet, target_col, value_row)
            worksheet.cell(row=target_row, column=target_col).value = round(amount, 2)
            worksheet.cell(row=target_row, column=target_col).number_format = "#,##0.00"
            result["written_count"] += 1
            log_callback(
                f"[엑셀입력] {image_path.name} -> {header_text}: {amount:.2f} (열{target_col}, 행{target_row})"
            )
        except Exception as exc:
            result["error_count"] += 1
            log_callback(f"[오류] {image_path.name}: {exc}")
        finally:
            if progress_callback:
                progress_callback(index, total_images, image_path.name)

    try:
        workbook.save(str(excel_path))
    except PermissionError as exc:
        raise RuntimeError("Excel file is open and cannot be saved. Close it and retry.") from exc

    log_callback(f"[엑셀 완료] 저장: {excel_path}")
    return result

def write_amounts_to_excel(excel_path, category_amounts, uncategorized_amounts=None, log_callback=print):
    if not OPENPYXL_AVAILABLE:
        raise RuntimeError("openpyxl 모듈이 설치되어 있지 않습니다.")

    excel_path = Path(excel_path)
    if not excel_path.exists():
        raise FileNotFoundError(f"엑셀 파일을 찾을 수 없습니다: {excel_path}")

    keep_vba = excel_path.suffix.lower() == ".xlsm"
    workbook = load_workbook(str(excel_path), keep_vba=keep_vba)
    worksheet = get_sum_worksheet(workbook)

    if uncategorized_amounts is None:
        uncategorized_amounts = []

    amounts_to_write = {}
    for category, values in category_amounts.items():
        valid_values = [round(abs(float(v)), 2) for v in values if v is not None]
        if valid_values:
            amounts_to_write[category] = valid_values
    unc_values = [round(abs(float(v)), 2) for v in uncategorized_amounts if v is not None]
    if unc_values:
        amounts_to_write[UNCATEGORIZED_FOLDER_NAME] = unc_values

    if not amounts_to_write:
        log_callback("[엑셀] 입력할 금액이 없어 저장만 진행합니다.")
        workbook.save(str(excel_path))
        return {
            "excel_path": str(excel_path),
            "sheet_name": worksheet.title,
            "written": {},
            "unmapped": {},
            "header_row": None,
            "value_row": None,
            "written_count": 0,
        }

    header_row, header_map = detect_sum_header_row_and_columns(worksheet, list(amounts_to_write.keys()))
    value_row = header_row + 1

    written = {}
    unmapped = {}
    written_count = 0

    for category, amount_list in amounts_to_write.items():
        target_col = header_map.get(category)
        if target_col is None:
            unmapped[category] = amount_list
            log_callback(f"[엑셀 미매핑] {category}: sum 시트 헤더를 찾지 못했습니다.")
            continue

        row = find_next_empty_row(worksheet, target_col, value_row)
        written_rows = []
        for amount in amount_list:
            worksheet.cell(row=row, column=target_col).value = round(amount, 2)
            worksheet.cell(row=row, column=target_col).number_format = "#,##0.00"
            written_rows.append({"row": row, "value": round(amount, 2)})
            log_callback(f"[엑셀 입력] {category}: {amount:.2f} (열 {target_col}, 행 {row})")
            written_count += 1
            row += 1
        written[category] = written_rows

    try:
        workbook.save(str(excel_path))
    except PermissionError as exc:
        raise RuntimeError("엑셀 파일이 열려 있어 저장할 수 없습니다. 파일을 닫고 다시 실행해 주세요.") from exc

    log_callback(f"[엑셀 완료] 저장: {excel_path}")
    return {
        "excel_path": str(excel_path),
        "sheet_name": worksheet.title,
        "written": written,
        "unmapped": unmapped,
        "header_row": header_row,
        "value_row": value_row,
        "written_count": written_count,
    }


def copy_excel_with_folder_name(excel_path, folder_path, log_callback=print):
    excel_path = Path(excel_path)
    folder_path = Path(folder_path)
    if not excel_path.exists():
        raise FileNotFoundError(f"복사할 엑셀 파일이 없습니다: {excel_path}")
    if not folder_path.exists():
        raise FileNotFoundError(f"선택 폴더가 없습니다: {folder_path}")

    destination = folder_path / f"{folder_path.name}{excel_path.suffix}"
    if destination.resolve() == excel_path.resolve():
        log_callback(f"[엑셀 복사] 원본과 대상이 동일하여 복사를 생략합니다: {destination}")
        return str(destination)

    shutil.copy2(str(excel_path), str(destination))
    log_callback(f"[엑셀 복사 완료] {destination}")
    return str(destination)


class ExpenseAutoApp:
    def __init__(self):
        self.root = Tk()
        self.root.title(APP_NAME)
        self.root.geometry("1140x700")
        self.root.minsize(1020, 620)

        self.selected_folder = StringVar(value=DEFAULT_INPUT_FOLDER)
        self.selected_work_folder2 = StringVar(value=DEFAULT_INPUT_FOLDER)
        self.selected_xlsx = StringVar()
        self.run_mode = StringVar(value="1")
        self.mode_buttons = []
        self.is_running = False
        self.log_queue = queue.Queue()
        self._active_processing_log = None
        self.desc_path = get_app_dir() / DESC_FILE_NAME

        self.help_window = None
        self.help_text_box = None

        self._icon_image = None
        self._preview_photo = None
        self._preview_source_path = None
        self._background_photo = None
        self._background_source = self._load_background_source()
        self._background_resize_job = None
        self._last_background_size = None

        self._setup_app_icon()
        self._build_menu()
        self._build_ui()

        self.root.bind("<F1>", self.open_help)
        self.root.bind("<Configure>", self._on_root_resize)
        self._update_background(force=True)
        self._update_progress_ui(0, 0, None)

        self.root.after(100, self._drain_queue)

    def _build_menu(self):
        menu_bar = tk.Menu(self.root)
        help_menu = tk.Menu(menu_bar, tearoff=0)
        help_menu.add_command(label="도움말(F1)", command=self.open_help)
        menu_bar.add_cascade(label="도움말", menu=help_menu)
        self.root.configure(menu=menu_bar)

    def _setup_app_icon(self):
        icon_path = find_existing_file(APP_ICON_CANDIDATES)
        if icon_path is None:
            return

        try:
            if icon_path.suffix.lower() == ".ico":
                self.root.iconbitmap(default=str(icon_path))
        except Exception:
            pass

        try:
            with Image.open(icon_path) as icon_img:
                icon_img = ImageOps.exif_transpose(icon_img).convert("RGBA")
                icon_img.thumbnail((256, 256), RESAMPLE_LANCZOS)
                self._icon_image = ImageTk.PhotoImage(icon_img)
            self.root.iconphoto(True, self._icon_image)
        except Exception:
            return

    def _load_background_source(self):
        bg_path = find_existing_file(BACKGROUND_IMAGE_CANDIDATES)
        if bg_path is None:
            return None
        try:
            with Image.open(bg_path) as bg_img:
                return ImageOps.exif_transpose(bg_img).convert("RGB")
        except Exception:
            return None

    def _build_ui(self):
        self.background_label = tk.Label(self.root, bd=0, highlightthickness=0)
        self.background_label.place(x=0, y=0, relwidth=1, relheight=1)
        self.background_label.lower()

        wrapper = ttk.Frame(self.root, padding=14)
        wrapper.place(x=0, y=0, relwidth=1, relheight=1)

        ttk.Label(wrapper, text=APP_NAME, font=("Malgun Gothic", 15, "bold")).pack(anchor="w", pady=(0, 10))

        folder_frame = ttk.Frame(wrapper)
        folder_frame.pack(fill="x", pady=(0, 10))
        ttk.Label(folder_frame, text="대상 폴더").pack(side="left")

        self.browse_button = ttk.Button(folder_frame, text="찾아보기", command=self.choose_folder)
        self.browse_button.pack(side="left", padx=(8, 8))

        self.folder_entry = ttk.Entry(folder_frame, textvariable=self.selected_folder, width=56)
        self.folder_entry.pack(side="left", fill="x", expand=True)

        feature_frame = ttk.Frame(wrapper)
        feature_frame.pack(fill="x", pady=(0, 10))

        self.mode_frame = ttk.LabelFrame(feature_frame, text="기능선택1", padding=8)
        self.mode_frame.pack(side="left", fill="both", expand=True, padx=(0, 6))
        mode_inner = ttk.Frame(self.mode_frame)
        mode_inner.pack(fill="x")

        mode_left = ttk.Frame(mode_inner)
        mode_left.pack(side="left", fill="x", expand=True)

        mode_right = ttk.Frame(mode_inner)
        mode_right.pack(side="right", anchor="n", padx=(12, 0))

        for mode_code, (label, _, _) in RUN_MODE_OPTIONS.items():
            rb = ttk.Radiobutton(mode_left, text=label, variable=self.run_mode, value=mode_code)
            rb.pack(anchor="w")
            self.mode_buttons.append(rb)

        self.run_button = ttk.Button(mode_right, text="실행", width=12, command=self.start_processing)
        self.run_button.pack(fill="x")

        self.quit_button = ttk.Button(mode_right, text="종료", width=12, command=self.root.destroy)
        self.quit_button.pack(fill="x", pady=(6, 0))

        self.status_label = ttk.Label(mode_right, text="대기 중")
        self.status_label.pack(anchor="e", pady=(8, 0))

        self.mode2_frame = ttk.LabelFrame(feature_frame, text="기능선택2", padding=8)
        self.mode2_frame.pack(side="left", fill="both", expand=True, padx=(6, 0))

        mode2_inner = ttk.Frame(self.mode2_frame)
        mode2_inner.pack(fill="x")

        row0 = ttk.Frame(mode2_inner)
        row0.pack(fill="x", pady=(0, 6))
        self.work_folder2_button = ttk.Button(row0, text="작업폴더", width=12, command=self.choose_work_folder2)
        self.work_folder2_button.pack(side="left")
        self.work_folder2_entry = ttk.Entry(row0, textvariable=self.selected_work_folder2, width=48)
        self.work_folder2_entry.pack(side="left", fill="x", expand=True, padx=(8, 0))

        row1 = ttk.Frame(mode2_inner)
        row1.pack(fill="x", pady=(0, 6))
        self.xlsx_open_button = ttk.Button(row1, text="xlsx 열기", width=12, command=self.choose_xlsx_file)
        self.xlsx_open_button.pack(side="left")
        self.xlsx_entry = ttk.Entry(row1, textvariable=self.selected_xlsx, width=48)
        self.xlsx_entry.pack(side="left", fill="x", expand=True, padx=(8, 0))

        row2 = ttk.Frame(mode2_inner)
        row2.pack(fill="x")
        self.run_match_button = ttk.Button(row2, text="실행", width=12, command=self.start_pptx_xlsx_processing)
        self.run_match_button.pack(side="left")

        split_frame = ttk.Panedwindow(wrapper, orient="horizontal")
        split_frame.pack(fill="both", expand=True, pady=(0, 0))

        left_frame = ttk.Frame(split_frame, padding=(0, 0, 8, 0))
        right_frame = ttk.Frame(split_frame, padding=(8, 0, 0, 0))
        split_frame.add(left_frame, weight=3)
        split_frame.add(right_frame, weight=2)

        progress_frame = ttk.LabelFrame(left_frame, text="진행 상태", padding=10)
        progress_frame.pack(fill="x", pady=(0, 8))

        self.progress_bar = ttk.Progressbar(progress_frame, mode="determinate")
        self.progress_bar.pack(fill="x")

        self.progress_detail_label = ttk.Label(progress_frame, text="0 / 0")
        self.progress_detail_label.pack(anchor="w", pady=(6, 0))

        self.current_file_label = ttk.Label(progress_frame, text="현재 파일: -")
        self.current_file_label.pack(anchor="w", pady=(2, 0))

        log_frame = ttk.LabelFrame(left_frame, text="처리 로그", padding=6)
        log_frame.pack(fill="both", expand=True)
        self.log_box = ScrolledText(log_frame, height=22, font=("Consolas", 10))
        self.log_box.pack(fill="both", expand=True)
        self.log_box.tag_configure("processing_highlight", foreground="#d97706", font=("Malgun Gothic", 10, "bold"))
        self.log_box.tag_configure("summary_highlight", foreground="#d97706", font=("Malgun Gothic", 10, "bold"))
        self.log_box.configure(state="disabled")

        preview_frame = ttk.LabelFrame(right_frame, text="인식 중인 이미지", padding=8)
        preview_frame.pack(fill="both", expand=True)

        self.preview_image_label = tk.Label(
            preview_frame,
            text="처리 중인 이미지가 여기에 표시됩니다.",
            anchor="center",
            justify="center",
            bg=PREVIEW_BG_COLOR,
            fg="#334155",
            relief="groove",
            bd=1,
        )
        self.preview_image_label.pack(fill="both", expand=True)
        self.preview_image_label.bind("<Configure>", self._on_preview_resize)

        self.preview_name_label = ttk.Label(preview_frame, text="파일: -")
        self.preview_name_label.pack(anchor="w", pady=(8, 0))

    def choose_folder(self):
        folder = filedialog.askdirectory(
            title="영수증 이미지 폴더 선택",
            initialdir=DEFAULT_INPUT_FOLDER,
        )
        if folder:
            self.selected_folder.set(folder)

    def choose_work_folder2(self):
        folder = filedialog.askdirectory(
            title="기능선택2 작업 폴더 선택",
            initialdir=self._resolve_initial_dir(self.selected_work_folder2.get()),
        )
        if folder:
            self.selected_work_folder2.set(folder)

    def _resolve_initial_dir(self, selected_path):
        path = Path(str(selected_path or "").strip())
        if path.exists():
            return str(path.parent if path.is_file() else path)
        if Path(DEFAULT_INPUT_FOLDER).exists():
            return DEFAULT_INPUT_FOLDER
        return str(get_app_dir())

    def choose_xlsx_file(self):
        file_path = filedialog.askopenfilename(
            title="XLSX 파일 선택",
            initialdir=self._resolve_initial_dir(self.selected_xlsx.get() or self.selected_work_folder2.get()),
            filetypes=[("Excel files", "*.xlsx *.xlsm"), ("All files", "*.*")],
        )
        if file_path:
            self.selected_xlsx.set(file_path)

    def _on_root_resize(self, event):
        if event.widget is not self.root:
            return
        if self._background_resize_job:
            self.root.after_cancel(self._background_resize_job)
        self._background_resize_job = self.root.after(120, self._update_background)

    def _update_background(self, force=False):
        if self._background_source is None:
            self.background_label.configure(image="", bg="#edf2f7")
            return

        width = max(self.root.winfo_width(), 1)
        height = max(self.root.winfo_height(), 1)
        size = (width, height)
        if not force and size == self._last_background_size:
            return

        self._last_background_size = size
        bg_img = ImageOps.fit(self._background_source, size, method=RESAMPLE_LANCZOS)
        bg_img = bg_img.filter(ImageFilter.GaussianBlur(radius=BACKGROUND_BLUR_RADIUS))
        bg_img = Image.blend(bg_img, Image.new("RGB", size, "#ffffff"), 0.40)

        self._background_photo = ImageTk.PhotoImage(bg_img)
        self.background_label.configure(image=self._background_photo)
        self.background_label.lower()

    def _update_progress_ui(self, current, total, file_name):
        total_val = max(int(total or 0), 0)
        max_val = total_val if total_val > 0 else 1
        cur_val = min(max(int(current or 0), 0), max_val)
        self.progress_bar.configure(maximum=max_val, value=cur_val)
        self.progress_detail_label.configure(text=f"{cur_val} / {total_val}")
        if file_name:
            self.current_file_label.configure(text=f"현재 파일: {file_name}")
        elif total_val == 0:
            self.current_file_label.configure(text="현재 파일: -")

    def _clear_preview(self, text="처리 중인 이미지가 여기에 표시됩니다."):
        self._preview_source_path = None
        self._preview_photo = None
        self.preview_image_label.configure(image="", text=text)
        self.preview_name_label.configure(text="파일: -")

    def _on_preview_resize(self, _event):
        if self._preview_source_path:
            self._render_preview_image()

    def _render_preview_image(self):
        if not self._preview_source_path or not self._preview_source_path.exists():
            self._clear_preview("이미지 파일을 찾을 수 없습니다.")
            return

        box_w = max(self.preview_image_label.winfo_width() - 12, PREVIEW_MIN_SIZE[0])
        box_h = max(self.preview_image_label.winfo_height() - 12, PREVIEW_MIN_SIZE[1])

        try:
            with Image.open(self._preview_source_path) as img:
                img = ImageOps.exif_transpose(img).convert("RGB")
        except Exception as exc:
            self._clear_preview(f"미리보기 오류\n{exc}")
            return

        img.thumbnail((box_w, box_h), RESAMPLE_LANCZOS)
        canvas = Image.new("RGB", (box_w, box_h), PREVIEW_BG_COLOR)
        x_pos = (box_w - img.width) // 2
        y_pos = (box_h - img.height) // 2
        canvas.paste(img, (x_pos, y_pos))

        self._preview_photo = ImageTk.PhotoImage(canvas)
        self.preview_image_label.configure(image=self._preview_photo, text="")
        self.preview_name_label.configure(text=f"파일: {self._preview_source_path.name}")

    def show_preview_image(self, image_path):
        if not image_path:
            return
        self._preview_source_path = Path(image_path)
        self._render_preview_image()

    def _load_help_text(self):
        if not self.desc_path.exists():
            return None
        text = read_text_with_fallback(self.desc_path).strip()
        return text if text else "(desc.txt 파일이 비어 있습니다.)"

    def _show_help_window(self, help_text):
        if self.help_window is None or not self.help_window.winfo_exists():
            self.help_window = tk.Toplevel(self.root)
            self.help_window.title("도움말")
            self.help_window.geometry(HELP_WINDOW_GEOMETRY)
            self.help_window.minsize(560, 400)
            self.help_window.transient(self.root)
            self.help_window.protocol("WM_DELETE_WINDOW", self._close_help_window)

            self.help_text_box = ScrolledText(self.help_window, wrap="word", font=("Malgun Gothic", 10))
            self.help_text_box.pack(fill="both", expand=True, padx=12, pady=12)

        self.help_window.deiconify()
        self.help_window.lift()
        self.help_window.focus_force()

        self.help_text_box.configure(state="normal")
        self.help_text_box.delete("1.0", "end")
        self.help_text_box.insert(
            "1.0",
            f"[파일 위치] {self.desc_path}\n\n{help_text}",
        )
        self.help_text_box.configure(state="disabled")

    def _close_help_window(self):
        if self.help_window is not None and self.help_window.winfo_exists():
            self.help_window.destroy()
        self.help_window = None
        self.help_text_box = None

    def open_help(self, _event=None):
        help_text = self._load_help_text()
        if help_text is None:
            messagebox.showwarning(
                "도움말 파일 안내",
                f"프로그램 폴더에 '{DESC_FILE_NAME}' 파일이 없습니다.\n\n경로: {self.desc_path}",
            )
            return "break"
        self._show_help_window(help_text)
        return "break"

    def _activate_processing_log(self, message):
        if self._active_processing_log is not None:
            self._deactivate_processing_log(keep_history=True)
        self.log_box.insert("1.0", message + "\n", "processing_highlight")
        self._active_processing_log = message

    def _deactivate_processing_log(self, keep_history=True):
        if self._active_processing_log is None:
            return
        prev_message = self._active_processing_log
        self.log_box.delete("1.0", "2.0")
        if keep_history:
            self.log_box.insert("end", prev_message + "\n")
        self._active_processing_log = None

    def append_log(self, message, style=None):
        self.log_box.configure(state="normal")
        display_message = message
        is_processing = False

        if isinstance(message, str) and message.startswith(LOG_PROCESSING_MARKER):
            display_message = message[len(LOG_PROCESSING_MARKER) :]
            is_processing = True

        if is_processing:
            self._activate_processing_log(display_message)
        elif style == "summary_highlight":
            if self._active_processing_log is not None:
                self._deactivate_processing_log(keep_history=True)
            self.log_box.insert("end", display_message + "\n", "summary_highlight")
        else:
            if self._active_processing_log is not None:
                self._deactivate_processing_log(keep_history=True)
            self.log_box.insert("end", display_message + "\n")
        self.log_box.see("end")
        self.log_box.configure(state="disabled")

    def set_running_state(self, running):
        self.is_running = running
        state = "disabled" if running else "normal"
        self.run_button.configure(state=state)
        self.browse_button.configure(state=state)
        self.folder_entry.configure(state=state)
        self.work_folder2_button.configure(state=state)
        self.work_folder2_entry.configure(state=state)
        self.xlsx_open_button.configure(state=state)
        self.xlsx_entry.configure(state=state)
        self.run_match_button.configure(state=state)
        for rb in self.mode_buttons:
            rb.configure(state=state)
        self.status_label.configure(text="처리 중..." if running else "대기 중")

    def start_processing(self):
        if self.is_running:
            return

        folder = self.selected_folder.get().strip()
        if not folder:
            messagebox.showwarning("입력 필요", "먼저 대상 폴더를 선택해 주세요.")
            return
        if not os.path.isdir(folder):
            messagebox.showerror("폴더 오류", "선택한 경로가 올바른 폴더가 아닙니다.")
            return

        mode_code = self.run_mode.get()
        mode_label, do_ocr, make_ppt = resolve_run_mode(mode_code)

        if make_ppt and not PPTX_AVAILABLE:
            messagebox.showerror("모듈 오류", "python-pptx 모듈이 없어 PPT를 만들 수 없습니다.")
            return

        if do_ocr and not os.path.exists(pytesseract.pytesseract.tesseract_cmd):
            messagebox.showerror(
                "Tesseract 경로 오류",
                f"Tesseract 실행파일을 찾을 수 없습니다.\n\n경로: {pytesseract.pytesseract.tesseract_cmd}",
            )
            return

        self.log_box.configure(state="normal")
        self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")
        self._active_processing_log = None
        self._update_progress_ui(0, 0, None)
        self._clear_preview()

        self.append_log(f"[시작] 폴더: {folder}")
        self.append_log(f"[모드] {mode_label}")

        self.set_running_state(True)

        worker = threading.Thread(
            target=self._process_worker,
            args=(folder, do_ocr, make_ppt),
            daemon=True,
        )
        worker.start()

    def start_pptx_xlsx_processing(self):
        if self.is_running:
            return

        work_folder = self.selected_work_folder2.get().strip()
        xlsx_path = self.selected_xlsx.get().strip()

        if not work_folder:
            messagebox.showwarning("입력 필요", "작업폴더를 선택해 주세요.")
            return
        if not xlsx_path:
            messagebox.showwarning("입력 필요", "xlsx 파일을 선택해 주세요.")
            return
        if not os.path.isdir(work_folder):
            messagebox.showerror("폴더 오류", f"작업폴더를 찾을 수 없습니다.\n\n{work_folder}")
            return
        if not os.path.isfile(xlsx_path):
            messagebox.showerror("파일 오류", f"xlsx 파일을 찾을 수 없습니다.\n\n{xlsx_path}")
            return
        if Path(xlsx_path).suffix.lower() not in EXCEL_EXTENSIONS:
            messagebox.showerror("파일 형식 오류", "xlsx/xlsm 파일만 선택할 수 있습니다.")
            return
        if not OPENPYXL_AVAILABLE:
            messagebox.showerror("모듈 오류", "openpyxl 모듈이 없어 XLSX를 처리할 수 없습니다.")
            return
        if not os.path.exists(pytesseract.pytesseract.tesseract_cmd):
            messagebox.showerror(
                "Tesseract 경로 오류",
                f"Tesseract 실행파일을 찾을 수 없습니다.\n\n경로: {pytesseract.pytesseract.tesseract_cmd}",
            )
            return

        self.log_box.configure(state="normal")
        self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")
        self._active_processing_log = None
        self._update_progress_ui(0, 0, None)
        self._clear_preview("기능선택2 실행 중에는 이미지 미리보기를 사용하지 않습니다.")

        self.append_log(f"[시작] 기능선택2 작업폴더: {work_folder}")
        self.append_log(f"[시작] 기능선택2 XLSX: {xlsx_path}")

        self.set_running_state(True)

        worker = threading.Thread(
            target=self._process_pptx_xlsx_worker,
            args=(work_folder, xlsx_path),
            daemon=True,
        )
        worker.start()

    def _process_pptx_xlsx_worker(self, work_folder, xlsx_path):
        try:
            result = write_folder_images_to_excel_by_headers(
                work_folder,
                xlsx_path,
                log_callback=lambda msg: self.log_queue.put(("log", msg)),
                progress_callback=lambda cur, total, name: self.log_queue.put(("progress", (cur, total, name))),
            )
            self.log_queue.put(("done_pptx_xlsx", result))
        except Exception as exc:
            self.log_queue.put(("fatal", str(exc)))

    def _process_worker(self, folder, do_ocr, make_ppt):
        try:
            result = {
                "classification": None,
                "ppt": None,
                "folder": folder,
            }

            if do_ocr:
                result["classification"] = process_images(
                    folder,
                    log_callback=lambda msg: self.log_queue.put(("log", msg)),
                    progress_callback=lambda cur, total, name: self.log_queue.put(("progress", (cur, total, name))),
                    preview_callback=lambda path: self.log_queue.put(("preview", path)),
                )
            else:
                self.log_queue.put(("progress", (0, 0, None)))
                self.log_queue.put(("preview", None))

            if make_ppt:
                self.log_queue.put(("log", "[PPT] 생성 시작"))
                result["ppt"] = create_ppt_from_subfolders(
                    folder,
                    log_callback=lambda msg: self.log_queue.put(("log", msg)),
                )

            self.log_queue.put(("done", result))
        except Exception as exc:
            self.log_queue.put(("fatal", str(exc)))

    def _drain_queue(self):
        while not self.log_queue.empty():
            kind, payload = self.log_queue.get_nowait()

            if kind == "log":
                self.append_log(payload)
                continue

            if kind == "progress":
                cur, total, file_name = payload
                self._update_progress_ui(cur, total, file_name)
                continue

            if kind == "preview":
                self.show_preview_image(payload)
                continue

            if kind == "done_pptx_xlsx":
                self.append_log("[완료] 기능선택2 작업이 종료되었습니다.")
                self.set_running_state(False)

                total_images = int(payload.get("total_images", 0))
                error_count = int(payload.get("error_count", 0))
                written_count = int(payload.get("written_count", 0))
                self._update_progress_ui(total_images, total_images, None)

                self.append_log(
                    f"[처리결과] 전체 {total_images}건 / 오류 {error_count}건 / 입력 {written_count}건",
                    style="summary_highlight",
                )

                unmatched_count = int(payload.get("unmatched_count", 0))
                no_amount_count = int(payload.get("no_amount_count", 0))
                self.append_log(
                    f"[결과상세] 헤더미일치 {unmatched_count}건 / 금액없음 {no_amount_count}건",
                    style="summary_highlight",
                )
                self.append_log(f"[엑셀] 저장: {payload.get('excel_path', '-')}")
                messagebox.showinfo("안내", "기능선택2 작업이 완료되었습니다.")
                continue

            if kind == "done":
                self.append_log("[완료] 작업이 종료되었습니다.")
                self.set_running_state(False)

                cls = payload.get("classification")
                ppt = payload.get("ppt")
                total_count = 0
                final_error_count = 0
                result_rows = {}
                if cls:
                    self._update_progress_ui(cls["total"], cls["total"], None)
                    total_count = int(cls.get("total", 0))
                    final_error_count = int(cls.get("errors", 0))
                    result_rows = cls.get("folder_counts", {})
                elif ppt:
                    total_count = int(ppt.get("images_total", 0))
                    result_rows = ppt.get("folder_image_counts", {})

                if result_rows:
                    for folder_name, processed_count in result_rows.items():
                        self.append_log(
                            f"[처리결과] 전체 {total_count}건 / 오류 {final_error_count}건 / {folder_name} {processed_count}건",
                            style="summary_highlight",
                        )
                elif cls or ppt:
                    self.append_log(
                        f"[처리결과] 전체 {total_count}건 / 오류 {final_error_count}건 / 처리폴더 없음 0건",
                        style="summary_highlight",
                    )

                if ppt:
                    self.append_log(
                        f"[PPT] 폴더 {ppt['folders_used']}개 / 이미지 {ppt['images_total']}장 / 슬라이드 {ppt['slides_total']}장"
                    )
                    self.append_log(f"[PPT] 저장: {ppt['output_path']}")

                if cls is None and ppt is None:
                    self.append_log("[처리결과] 처리된 항목이 없습니다.", style="summary_highlight")

                messagebox.showinfo("안내", "미분류는 수동으로 처리 필요합니다")
                continue

            if kind == "fatal":
                self.append_log(f"[치명적 오류] {payload}")
                self.set_running_state(False)
                messagebox.showerror("오류", payload)

        self.root.after(100, self._drain_queue)

    def run(self):
        self.root.mainloop()


def main():
    app = ExpenseAutoApp()
    app.run()


if __name__ == "__main__":
    main()

